from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()

# Define a function to add title & content slides
def add_content_slide(title, content):
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = content

# Define a function to add a slide with a colored shape for visual appeal
def add_styled_slide(title, content, color=(91, 155, 213)):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title

    # Add a colored rectangle for style
    left = Inches(0)
    top = Inches(5.5)
    width = Inches(10)
    height = Inches(1)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*color)
    line = shape.line
    line.color.rgb = RGBColor(255, 255, 255)

    # Add content
    slide.placeholders[1].text = content

# Title Slide
slide_layout = prs.slide_layouts[0]  # Title Slide
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "Case Study: Jeevanksh"
slide.placeholders[1].text = "Comprehensive Analysis & Solutions for Supply Chain & Market Access"

# Add slides
add_styled_slide("Why This Product/Service?",
"""
- Focus on organic, traceable produce
- Empower smallholder farmers
- Meet growing health-conscious consumer demand
- Aligns with sustainability and ESG goals
""")

add_styled_slide("Market Strategy",
"""
- Direct-to-Consumer (D2C) online marketplace
- Partnerships with local cooperatives & SHGs
- Digital storytelling: farmer stories & process videos
- Premium branding for trust and loyalty
""")

add_styled_slide("Competitive Landscape",
"""
Main Competitors:
1. Organic India - National brand
2. 24 Mantra Organic - Retail chain supply
3. BigBasket Fresho - Convenience
4. Local Farmer Markets - Local trust
5. Amazon Pantry - Scale & logistics
""")

add_styled_slide("Comparison & New Strategies",
"""
- Blockchain-based traceability system
- Farmer digital ID & QR codes for transparency
- Subscription boxes for loyal consumers
- B2B supply for hotels, restaurants, organic cafes
""")

add_styled_slide("Value Proposition",
"""
- Verified, authentic organic produce
- Fair income and empowerment for farmers
- Healthy, chemical-free food for consumers
- Contribute to sustainable agriculture
""")

add_styled_slide("Revenue Model",
"""
- Direct online sales via website & app
- Monthly subscription plans
- B2B bulk supply to businesses
- Premium pricing for traceable authenticity
""")

add_styled_slide("Risks and Challenges",
"""
- Ensuring authenticity at scale
- Farmer education and compliance
- Managing perishable goods efficiently
- Competing with retail discounts and large players
""")

add_styled_slide("Conclusions",
"""
- Leverage tech for supply chain transparency
- Build strong farmer relationships & training programs
- Focus on brand storytelling & community
- Innovate with loyalty programs & B2B partnerships
""")

# Problem 1 Solution Slide
add_styled_slide("Problem 1: Enhancing Supply Chain Transparency",
"""
Framework:
✔ Blockchain records for each step (farm → consumer)
✔ QR codes on every product: farmer details, harvest date, testing reports
✔ Third-party audits & certifications
✔ Farmer app for real-time updates & compliance tracking
""", color=(112, 173, 71))  # Green accent

# Problem 2 Solution Slide
add_styled_slide("Problem 2: Expanding Market Access",
"""
Solutions:
✔ Local aggregation centers for efficient logistics
✔ Digital literacy training & market awareness for farmers
✔ Regional e-marketplaces connecting farmers directly to cities
✔ Partnership with B2B buyers: restaurants, cafes, hotels
✔ Brand ambassadors to promote local produce
""", color=(237, 125, 49))  # Orange accent

# Final Thank You Slide
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "Thank You"
slide.placeholders[1].text = "Website: https://www.jeevanksh.com\nFor more information, connect with the Jeevanksh team."

# Save the presentation
prs.save("Jeevanksh_Case_Study_Complete.pptx")
print("✅ Presentation saved as 'Jeevanksh_Case_Study_Complete.pptx'")
